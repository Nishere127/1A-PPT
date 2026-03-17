"""
将多张图片导出为 PPT：一页一图，可选标题/备注。
"""

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def ensure_blank_template(base_dir: Path | None = None) -> str:
    """若 templates/blank.pptx 不存在则创建 16:9 空白模板。返回模板路径。"""
    base = base_dir or Path(__file__).resolve().parent.parent
    tpl_dir = base / "templates"
    tpl_dir.mkdir(parents=True, exist_ok=True)
    tpl_path = tpl_dir / "blank.pptx"
    if tpl_path.exists():
        return str(tpl_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.save(str(tpl_path))
    return str(tpl_path)


def export_images_to_ppt(image_items: list[dict], output_path: str, template_path: str | None = None, base_dir: Path | None = None) -> str:
    """
    将多张图片导出为 PPT，一页一图。可选每页标题（来自 prompt/title）。
    image_items: 每项至少含 "bytes"（bytes）或 "path"（str）；可选 "prompt" 或 "title" 作为页标题。
    """
    base = base_dir or Path(__file__).resolve().parent.parent
    tpl = template_path or ensure_blank_template(base)
    prs = Presentation(tpl)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layouts = prs.slide_layouts
    blank = layouts[6] if len(layouts) > 6 else layouts[-1]
    img_w, img_h = Inches(10), Inches(5.5)
    left_img = Inches(1.65)
    top_img = Inches(0.8)
    title_h = Inches(0.6)
    for item in image_items:
        slide = prs.slides.add_slide(blank)
        title_text = item.get("title") or item.get("prompt") or ""
        if title_text and len(title_text) > 80:
            title_text = title_text[:77] + "..."
        if title_text:
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), title_h)
            p = tx.text_frame.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(14)
            p.font.bold = True
        if "bytes" in item:
            stream = BytesIO(item["bytes"])
            slide.shapes.add_picture(stream, left_img, top_img, width=img_w, height=img_h)
        elif "path" in item:
            path = Path(item["path"])
            if path.exists():
                slide.shapes.add_picture(str(path), left_img, top_img, width=img_w, height=img_h)
    prs.save(output_path)
    return output_path
