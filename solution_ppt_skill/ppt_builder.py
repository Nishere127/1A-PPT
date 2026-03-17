"""
根据大纲生成 PPT 文件。支持用户先编辑大纲再生成。
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def ensure_blank_template(base_dir: Path | None = None) -> str:
    """若 templates/blank.pptx 不存在则创建 16:9 空白模板。返回模板路径。"""
    base = base_dir or Path(__file__).resolve().parent.parent
    tpl_dir = base / "templates"
    tpl_dir.mkdir(parents=True, exist_ok=True)
    tpl_path = tpl_dir / "blank.pptx"
    if tpl_path.exists():
        return str(tpl_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.save(str(tpl_path))
    return str(tpl_path)


def create_ppt_from_outline(outline: list, template_path: str, output_path: str) -> str:
    """
    根据大纲和空白模板生成 PPT。outline: [{"title": "...", "points": ["...", ...]}, ...]
    """
    prs = Presentation(template_path)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layouts = prs.slide_layouts
    blank = layouts[6] if len(layouts) > 6 else layouts[-1]
    for item in outline:
        title = item.get("title") or "无标题"
        points = item.get("points") or []
        slide = prs.slides.add_slide(blank)
        left, top, w, h = Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        tx = slide.shapes.add_textbox(left, top, w, h)
        p = tx.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        body_top = Inches(1.5)
        for i, pt in enumerate(points[:5]):
            bx = slide.shapes.add_textbox(Inches(0.5), body_top + Inches(i * 0.9), Inches(12.333), Inches(0.7))
            bp = bx.text_frame.paragraphs[0]
            bp.text = f"• {pt}"
            bp.font.size = Pt(14)
    prs.save(output_path)
    return output_path
